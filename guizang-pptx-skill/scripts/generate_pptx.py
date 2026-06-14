#!/usr/bin/env python3
"""
guizang-pptx-skill — PPTX generation engine
Maps guizang-ppt-skill's design system to native PowerPoint (.pptx) output.

Usage:
  python generate_pptx.py spec.json output.pptx
  python generate_pptx.py -  output.pptx  < spec.json   # stdin

Spec format: see build_spec() or README.
"""

import json
import sys
import os
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


# ═══════════════════════════════════════════════
# THEME DEFINITIONS (from guizang-ppt-skill)
# ═══════════════════════════════════════════════

STYLE_A_THEMES = {
    "ink-classic": {
        "name": "墨水经典",
        "ink": "#0a0a0b", "ink_rgb": (10, 10, 11),
        "paper": "#f1efea", "paper_rgb": (241, 239, 234),
        "paper_tint": "#e8e5de", "ink_tint": "#18181a",
        "accent": "#0a0a0b",  # same as ink for A style
    },
    "indigo-porcelain": {
        "name": "靛蓝瓷",
        "ink": "#0a1f3d", "ink_rgb": (10, 31, 61),
        "paper": "#f1f3f5", "paper_rgb": (241, 243, 245),
        "paper_tint": "#e4e8ec", "ink_tint": "#152a4a",
        "accent": "#0a1f3d",
    },
    "forest-ink": {
        "name": "森林墨",
        "ink": "#1a2e1f", "ink_rgb": (26, 46, 31),
        "paper": "#f5f1e8", "paper_rgb": (245, 241, 232),
        "paper_tint": "#ece7da", "ink_tint": "#253d2c",
        "accent": "#1a2e1f",
    },
    "kraft-paper": {
        "name": "牛皮纸",
        "ink": "#2a1e13", "ink_rgb": (42, 30, 19),
        "paper": "#eedfc7", "paper_rgb": (238, 223, 199),
        "paper_tint": "#e0d0b6", "ink_tint": "#3a2a1d",
        "accent": "#2a1e13",
    },
    "dune": {
        "name": "沙丘",
        "ink": "#1f1a14", "ink_rgb": (31, 26, 20),
        "paper": "#f0e6d2", "paper_rgb": (240, 230, 210),
        "paper_tint": "#e3d7bf", "ink_tint": "#2d2620",
        "accent": "#1f1a14",
    },
}

STYLE_B_THEMES = {
    "ikb": {
        "name": "克莱因蓝 IKB",
        "paper": "#fafaf8", "paper_rgb": (250, 250, 248),
        "ink": "#0a0a0a", "ink_rgb": (10, 10, 10),
        "grey_1": "#f0f0ee", "grey_2": "#d4d4d2", "grey_3": "#737373",
        "accent": "#002FA7", "accent_rgb": (0, 47, 167),
        "accent_on": "#ffffff",
    },
    "lemon": {
        "name": "柠檬黄",
        "paper": "#fafaf8", "paper_rgb": (250, 250, 248),
        "ink": "#0a0a0a", "ink_rgb": (10, 10, 10),
        "grey_1": "#f0f0ee", "grey_2": "#d4d4d2", "grey_3": "#737373",
        "accent": "#FFD500", "accent_rgb": (255, 213, 0),
        "accent_on": "#0a0a0a",
    },
    "lemon-green": {
        "name": "柠檬绿",
        "paper": "#fafaf8", "paper_rgb": (250, 250, 248),
        "ink": "#0a0a0a", "ink_rgb": (10, 10, 10),
        "grey_1": "#f0f0ee", "grey_2": "#d4d4d2", "grey_3": "#737373",
        "accent": "#C5E803", "accent_rgb": (197, 232, 3),
        "accent_on": "#0a0a0a",
    },
    "safety-orange": {
        "name": "安全橙",
        "paper": "#fafaf8", "paper_rgb": (250, 250, 248),
        "ink": "#0a0a0a", "ink_rgb": (10, 10, 10),
        "grey_1": "#f0f0ee", "grey_2": "#d4d4d2", "grey_3": "#737373",
        "accent": "#FF6B35", "accent_rgb": (255, 107, 53),
        "accent_on": "#ffffff",
    },
}

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Font stacks
SERIF_ZH = "Noto Serif SC"
SERIF_EN = "Playfair Display"
SANS_ZH = "Microsoft YaHei UI"
SANS_EN = "Inter"
MONO = "Cascadia Code"


def hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def resolve_theme(spec: dict) -> dict:
    """Resolve theme from spec, supporting both style A and B."""
    style = spec.get("style", "A")
    theme_key = spec.get("theme", "ink-classic" if style == "A" else "ikb")

    if style == "B" or theme_key in STYLE_B_THEMES:
        return STYLE_B_THEMES.get(theme_key, STYLE_B_THEMES["ikb"])
    return STYLE_A_THEMES.get(theme_key, STYLE_A_THEMES["ink-classic"])


# ═══════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════

class PPTXBuilder:
    def __init__(self, spec: dict):
        self.spec = spec
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.theme = resolve_theme(spec)
        self.style = spec.get("style", "A")
        self.title = spec.get("title", "Presentation")
        self.author = spec.get("author", "")

    def _add_bg(self, slide, color_hex: str):
        """Set solid background color."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(color_hex)

    def _add_textbox(self, slide, left, top, width, height,
                     text: str, font_name: str = None, font_size: int = 18,
                     color: str = None, bold: bool = False, alignment=PP_ALIGN.LEFT,
                     line_spacing: float = 1.2):
        """Add a text box with single paragraph."""
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                          Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.alignment = alignment
        p.space_after = Pt(0)
        p.line_spacing = Pt(font_size * line_spacing)
        if color:
            p.font.color.rgb = hex_to_rgb(color)
        else:
            p.font.color.rgb = hex_to_rgb(self.theme["ink"])
        if font_name:
            p.font.name = font_name
        else:
            p.font.name = SANS_ZH if self.style == "B" else SERIF_ZH
        return txBox

    def _add_multiline_textbox(self, slide, left, top, width, height,
                                paragraphs: list, font_name: str = None):
        """Add a text box with multiple styled paragraphs.

        Each paragraph is a dict: {text, size, color, bold, alignment, spacing, font}
        """
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                          Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, para in enumerate(paragraphs):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = para.get("text", "")
            p.font.size = Pt(para.get("size", 18))
            p.font.bold = para.get("bold", False)
            p.alignment = para.get("alignment", PP_ALIGN.LEFT)
            p.space_after = Pt(para.get("space_after", 4))
            p.line_spacing = Pt(para.get("size", 18) * para.get("line_spacing", 1.3))

            clr = para.get("color", self.theme["ink"])
            p.font.color.rgb = hex_to_rgb(clr)

            fn = para.get("font", font_name)
            if fn:
                p.font.name = fn
            elif self.style == "B":
                p.font.name = SANS_EN
            else:
                p.font.name = para.get("font", font_name or SERIF_ZH)

        return txBox

    def _add_accent_bar(self, slide, left, top, width, height):
        """Add a solid accent-colored rectangle."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(self.theme["accent"])
        shape.line.fill.background()  # no outline
        return shape

    def _add_image(self, slide, left, top, width, height, image_path: str):
        """Add an image. Supports relative paths from spec dir."""
        if not os.path.isabs(image_path):
            base = self.spec.get("_base_dir", os.getcwd())
            image_path = os.path.join(base, image_path)
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(left), Inches(top),
                                      Inches(width), Inches(height))
        else:
            # Placeholder rectangle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(self.theme.get("grey_1", self.theme.get("paper_tint", "#e0e0e0")))
            shape.line.color.rgb = hex_to_rgb(self.theme.get("grey_2", "#cccccc"))
            shape.line.width = Pt(0.5)
            # Add placeholder text
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"[Image: {os.path.basename(image_path)}]"
            p.font.size = Pt(10)
            p.font.color.rgb = hex_to_rgb(self.theme.get("grey_3", "#999999"))
            p.alignment = PP_ALIGN.CENTER

    # ── Layout Builders ──

    def build_hero_cover(self, slide_spec: dict):
        """L1 / S01: Hero Cover — dark background + centered big title."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank
        is_dark = slide_spec.get("theme_class", "hero dark").endswith("dark")
        bg_color = self.theme["ink"] if is_dark else self.theme["paper"]
        text_color = self.theme["paper"] if is_dark else self.theme["ink"]
        self._add_bg(slide, bg_color)

        kicker = slide_spec.get("kicker", "")
        title_text = slide_spec.get("title", self.title)
        subtitle = slide_spec.get("subtitle", "")
        lead = slide_spec.get("lead", "")
        meta = slide_spec.get("meta", self.author)

        y = 1.8
        if kicker:
            self._add_textbox(slide, 1.5, y, 10, 0.5, kicker,
                              font_size=14, color=text_color,
                              font_name=SANS_ZH if self.style == "B" else MONO)
            y += 0.6

        # Main title — large, weight 200 for Swiss
        title_weight = 200 if self.style == "B" else 700
        tbox = self._add_textbox(slide, 1.5, y, 10, 1.8, title_text,
                                  font_size=52, color=text_color,
                                  font_name=SANS_EN if self.style == "B" else SERIF_ZH)
        y += 1.6

        if subtitle:
            self._add_textbox(slide, 1.5, y, 10, 0.8, subtitle,
                              font_size=28, color=text_color,
                              font_name=SANS_EN if self.style == "B" else SERIF_ZH)
            y += 0.9

        if lead:
            self._add_textbox(slide, 1.5, y, 8, 1.0, lead,
                              font_size=16, color=text_color,
                              font_name=SANS_ZH)
            y += 0.9

        if meta:
            self._add_textbox(slide, 1.5, y, 8, 0.6, meta,
                              font_size=12, color=text_color,
                              font_name=MONO)

        # Accent bar bottom-left
        self._add_accent_bar(slide, 1.5, 6.8, 1.5, 0.04)

    def build_act_divider(self, slide_spec: dict):
        """L2 / S03 Split: Chapter divider — minimal, bold."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        is_dark = slide_spec.get("theme_class", "hero light").endswith("dark")
        bg_color = self.theme["ink"] if is_dark else self.theme["paper"]
        text_color = self.theme["paper"] if is_dark else self.theme["ink"]
        self._add_bg(slide, bg_color)

        kicker = slide_spec.get("kicker", "")
        title_text = slide_spec.get("title", "")
        lead = slide_spec.get("lead", "")

        if self.style == "B":
            # Swiss: big accent block left side
            self._add_accent_bar(slide, 1.2, 2.0, 0.06, 3.5)

        y = 2.0
        if kicker:
            self._add_textbox(slide, 1.8, y, 10, 0.5, kicker,
                              font_size=13, color=text_color, font_name=MONO)
            y += 0.5

        title_fs = 48 if self.style == "B" else 54
        title_fn = SANS_EN if self.style == "B" else SERIF_ZH
        self._add_textbox(slide, 1.8, y, 10, 1.5, title_text,
                          font_size=title_fs, color=text_color, font_name=title_fn)
        y += 1.4

        if lead:
            self._add_textbox(slide, 1.8, y, 8, 0.8, lead,
                              font_size=16, color=text_color, font_name=SANS_ZH)

    def build_big_numbers(self, slide_spec: dict):
        """L3 / S06: Stat cards grid."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        is_dark = slide_spec.get("theme_class", "light").endswith("dark")
        bg_color = self.theme["ink"] if is_dark else self.theme["paper"]
        text_color = self.theme["paper"] if is_dark else self.theme["ink"]
        self._add_bg(slide, bg_color)

        title = slide_spec.get("title", "")
        kicker = slide_spec.get("kicker", "")
        lead = slide_spec.get("lead", "")
        stats = slide_spec.get("stats", [])

        y = 0.8
        if kicker:
            self._add_textbox(slide, 1.2, y, 10, 0.4, kicker,
                              font_size=12, color=text_color, font_name=MONO)
            y += 0.5
        if title:
            self._add_textbox(slide, 1.2, y, 10, 0.7, title,
                              font_size=30, color=text_color,
                              font_name=SANS_EN if self.style == "B" else SERIF_ZH)
            y += 0.8
        if lead:
            self._add_textbox(slide, 1.2, y, 10, 0.5, lead,
                              font_size=14, color=text_color, font_name=SANS_ZH)
            y += 0.6

        # Stat cards grid (3x2 or flexible)
        n = len(stats)
        if n == 0:
            return
        cols = min(n, 4)
        rows = (n + cols - 1) // cols
        card_w = 2.7
        card_h = 1.8
        gap_x = 0.25
        gap_y = 0.3
        start_x = 1.2
        start_y = max(y + 0.3, 2.5)

        for i, stat in enumerate(stats):
            r = i // cols
            c = i % cols
            cx = start_x + c * (card_w + gap_x)
            cy = start_y + r * (card_h + gap_y)

            # Card background
            if self.style == "B":
                card_bg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(cx), Inches(cy), Inches(card_w), Inches(card_h)
                )
                card_bg.fill.solid()
                card_bg.fill.fore_color.rgb = hex_to_rgb(self.theme["grey_1"])
                card_bg.line.fill.background()

            # Stat number
            nb = stat.get("number", stat.get("nb", ""))
            label = stat.get("label", "")
            note = stat.get("note", "")

            self._add_textbox(slide, cx + 0.2, cy + 0.15, card_w - 0.4, 0.3,
                              label, font_size=10,
                              color=self.theme.get("grey_3", text_color),
                              font_name=MONO)
            self._add_textbox(slide, cx + 0.2, cy + 0.45, card_w - 0.4, 0.8,
                              nb, font_size=36, bold=True,
                              color=self.theme["accent"] if self.style == "B" else text_color,
                              font_name=SANS_EN)
            self._add_textbox(slide, cx + 0.2, cy + 1.25, card_w - 0.4, 0.4,
                              note, font_size=10, color=text_color,
                              font_name=SANS_ZH)

    def build_quote_image(self, slide_spec: dict):
        """L4: Left text + right image."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        is_dark = slide_spec.get("theme_class", "light").endswith("dark")
        bg_color = self.theme["ink"] if is_dark else self.theme["paper"]
        text_color = self.theme["paper"] if is_dark else self.theme["ink"]
        self._add_bg(slide, bg_color)

        title = slide_spec.get("title", "")
        kicker = slide_spec.get("kicker", "")
        lead = slide_spec.get("lead", "")
        callout = slide_spec.get("callout", "")
        callout_src = slide_spec.get("callout_src", "")
        image = slide_spec.get("image", "")

        # Left column
        y = 1.2
        if kicker:
            self._add_textbox(slide, 1.0, y, 6, 0.4, kicker,
                              font_size=12, color=text_color, font_name=MONO)
            y += 0.5
        self._add_textbox(slide, 1.0, y, 6, 1.0, title,
                          font_size=32, color=text_color, font_name=SERIF_ZH if self.style == "A" else SANS_EN)
        y += 1.1
        if lead:
            self._add_textbox(slide, 1.0, y, 6, 1.2, lead,
                              font_size=14, color=text_color, font_name=SANS_ZH)
            y += 1.1

        # Callout at bottom of left column
        if callout:
            self._add_accent_bar(slide, 1.0, 5.8, 0.04, 0.8)
            self._add_textbox(slide, 1.2, 5.8, 6, 0.7, callout,
                              font_size=13, color=text_color,
                              font_name=SERIF_ZH if self.style == "A" else SANS_ZH)
            if callout_src:
                self._add_textbox(slide, 1.2, 6.4, 6, 0.4, callout_src,
                                  font_size=10, color=self.theme.get("grey_3", text_color),
                                  font_name=MONO)

        # Right column — image
        if image:
            self._add_image(slide, 7.5, 1.5, 5.0, 5.0, image)

    def build_image_grid(self, slide_spec: dict):
        """L5 / S15-S16: Image grid."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg_color = self.theme["paper"]
        self._add_bg(slide, bg_color)

        title = slide_spec.get("title", "")
        kicker = slide_spec.get("kicker", "")
        images = slide_spec.get("images", [])

        y = 0.8
        if kicker:
            self._add_textbox(slide, 1.0, y, 10, 0.4, kicker,
                              font_size=12, color=self.theme["ink"], font_name=MONO)
            y += 0.5
        if title:
            self._add_textbox(slide, 1.0, y, 10, 0.6, title,
                              font_size=26, color=self.theme["ink"],
                              font_name=SANS_EN if self.style == "B" else SERIF_ZH)
            y += 0.8

        n = len(images)
        if n == 0:
            return
        cols = min(n, 3)
        rows = (n + cols - 1) // cols
        img_w = 3.5
        img_h = 2.2
        gap = 0.3
        start_x = 1.0
        start_y = max(y + 0.2, 2.0)

        for i, img in enumerate(images):
            r = i // cols
            c = i % cols
            cx = start_x + c * (img_w + gap)
            cy = start_y + r * (img_h + gap + 0.35)
            if isinstance(img, str):
                self._add_image(slide, cx, cy, img_w, img_h, img)
            else:
                self._add_image(slide, cx, cy, img_w, img_h, img.get("path", ""))
                cap = img.get("caption", "")
                if cap:
                    self._add_textbox(slide, cx, cy + img_h + 0.05, img_w, 0.3,
                                      cap, font_size=10,
                                      color=self.theme.get("grey_3", "#737373"),
                                      font_name=SANS_ZH)

    def build_pipeline(self, slide_spec: dict):
        """L6 / S11: Horizontal flow steps."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg_color = self.theme["paper"]
        self._add_bg(slide, bg_color)

        title = slide_spec.get("title", "")
        kicker = slide_spec.get("kicker", "")
        pipelines = slide_spec.get("pipelines", [])

        y = 0.8
        if kicker:
            self._add_textbox(slide, 1.0, y, 10, 0.4, kicker,
                              font_size=12, color=self.theme["ink"], font_name=MONO)
            y += 0.5
        if title:
            self._add_textbox(slide, 1.0, y, 10, 0.6, title,
                              font_size=26, color=self.theme["ink"],
                              font_name=SANS_EN if self.style == "B" else SERIF_ZH)
            y += 0.8

        for pidx, pipe in enumerate(pipelines):
            label = pipe.get("label", "")
            steps = pipe.get("steps", [])
            if label:
                self._add_textbox(slide, 1.0, y, 10, 0.4, label,
                                  font_size=11, color=self.theme.get("grey_3", "#737373"),
                                  font_name=MONO)
                y += 0.5

            n = len(steps)
            step_w = 2.0
            step_gap = 0.35
            for i, step in enumerate(steps):
                sx = 1.0 + i * (step_w + step_gap)
                # Accent top line
                self._add_accent_bar(slide, sx, y, step_w, 0.03)
                self._add_textbox(slide, sx, y + 0.15, step_w, 0.35,
                                  f"{step.get('nb', i+1):02d}",
                                  font_size=24, bold=True,
                                  color=self.theme["accent"], font_name=SANS_EN)
                self._add_textbox(slide, sx, y + 0.55, step_w, 0.35,
                                  step.get("title", ""),
                                  font_size=14, bold=True,
                                  color=self.theme["ink"], font_name=SANS_EN)
                self._add_textbox(slide, sx, y + 0.9, step_w, 0.6,
                                  step.get("desc", ""),
                                  font_size=10, color=self.theme.get("grey_3", "#737373"),
                                  font_name=SANS_ZH)
            y += 1.7

    def build_hero_question(self, slide_spec: dict):
        """L7 / S09: Centered big question"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg_color = self.theme["ink"]
        self._add_bg(slide, bg_color)

        kicker = slide_spec.get("kicker", "")
        question = slide_spec.get("question", slide_spec.get("title", ""))
        lead = slide_spec.get("lead", "")

        y = 2.0
        if kicker:
            self._add_textbox(slide, 1.5, y, 10, 0.5, kicker,
                              font_size=13, color=self.theme["paper"],
                              font_name=MONO, alignment=PP_ALIGN.CENTER)
            y += 0.8

        # Big question text
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(10), Inches(2.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = question
        p.font.size = Pt(40)
        p.font.color.rgb = hex_to_rgb(self.theme["paper"])
        p.font.name = SERIF_ZH if self.style == "A" else SANS_EN
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = Pt(52)
        y += 2.5

        if lead:
            self._add_textbox(slide, 2, y, 9, 0.8, lead,
                              font_size=14, color=self.theme["paper"],
                              font_name=SANS_ZH, alignment=PP_ALIGN.CENTER)

    def build_big_quote(self, slide_spec: dict):
        """L8: Large quote slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        is_dark = slide_spec.get("theme_class", "dark").endswith("dark")
        bg_color = self.theme["ink"] if is_dark else self.theme["paper"]
        text_color = self.theme["paper"] if is_dark else self.theme["ink"]
        self._add_bg(slide, bg_color)

        kicker = slide_spec.get("kicker", "")
        quote = slide_spec.get("quote", "")
        source = slide_spec.get("source", "")
        source_meta = slide_spec.get("source_meta", "")

        y = 1.8
        if kicker:
            self._add_textbox(slide, 2, y, 9, 0.5, kicker,
                              font_size=12, color=text_color,
                              font_name=MONO, alignment=PP_ALIGN.CENTER)
            y += 0.8

        # Big quote
        txBox = slide.shapes.add_textbox(Inches(2), Inches(y), Inches(9), Inches(3.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f'"{quote}"' if not quote.startswith('"') else quote
        p.font.size = Pt(38)
        p.font.color.rgb = hex_to_rgb(text_color)
        p.font.name = SERIF_ZH if self.style == "A" else SANS_EN
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = Pt(50)
        y += 3.2

        if source:
            self._add_textbox(slide, 3, y + 0.5, 7, 0.5, f"— {source}",
                              font_size=14, color=text_color,
                              font_name=SANS_ZH, alignment=PP_ALIGN.CENTER)
        if source_meta:
            self._add_textbox(slide, 3, y + 1.0, 7, 0.4, source_meta,
                              font_size=11, color=text_color,
                              font_name=MONO, alignment=PP_ALIGN.CENTER)

    def build_compare(self, slide_spec: dict):
        """L9 / S08: Before/After comparison."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg_color = self.theme["paper"]
        self._add_bg(slide, bg_color)

        title = slide_spec.get("title", "")
        kicker = slide_spec.get("kicker", "")
        left = slide_spec.get("left", {})
        right = slide_spec.get("right", {})

        y = 0.8
        if kicker:
            self._add_textbox(slide, 1.0, y, 11, 0.4, kicker,
                              font_size=12, color=self.theme["ink"], font_name=MONO)
            y += 0.5
        if title:
            self._add_textbox(slide, 1.0, y, 11, 0.6, title,
                              font_size=26, color=self.theme["ink"],
                              font_name=SANS_EN if self.style == "B" else SERIF_ZH)
            y += 0.8

        # Two columns
        col_y = y + 0.3
        for ci, (col_data, is_left) in enumerate([(left, True), (right, False)]):
            cx = 1.0 if is_left else 6.8
            cw = 5.5
            opacity_color = self.theme.get("grey_1", "#f0f0ee")

            # Column bg
            col_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(cx), Inches(col_y), Inches(cw), Inches(5.0)
            )
            col_bg.fill.solid()
            col_bg.fill.fore_color.rgb = hex_to_rgb(opacity_color if is_left else self.theme["paper"])
            col_bg.line.color.rgb = hex_to_rgb(self.theme.get("grey_2", "#cccccc"))
            col_bg.line.width = Pt(0.5)

            # Border-left accent (thick for right/new)
            border_color = self.theme.get("grey_3", "#999999") if is_left else self.theme["accent"]
            self._add_accent_bar(slide, cx, col_y, 0.06, 5.0)

            ckicker = col_data.get("kicker", "")
            ctitle = col_data.get("title", "")
            items = col_data.get("items", [])

            cy = col_y + 0.3
            if ckicker:
                self._add_textbox(slide, cx + 0.3, cy, cw - 0.5, 0.35, ckicker,
                                  font_size=11, color=self.theme["ink"], font_name=MONO)
                cy += 0.45
            if ctitle:
                self._add_textbox(slide, cx + 0.3, cy, cw - 0.5, 0.5, ctitle,
                                  font_size=18, bold=True, color=self.theme["ink"],
                                  font_name=SANS_EN)
                cy += 0.55

            for item in items:
                self._add_textbox(slide, cx + 0.5, cy, cw - 0.8, 0.35,
                                  f"• {item}", font_size=12,
                                  color=self.theme["ink"], font_name=SANS_ZH)
                cy += 0.3

    def build_mixed_text_image(self, slide_spec: dict):
        """L10: Lead image + side text."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg_color = self.theme["paper"]
        self._add_bg(slide, bg_color)

        title = slide_spec.get("title", "")
        kicker = slide_spec.get("kicker", "")
        body = slide_spec.get("body", "")
        callout = slide_spec.get("callout", "")
        callout_src = slide_spec.get("callout_src", "")
        image = slide_spec.get("image", "")

        y = 0.8
        if kicker:
            self._add_textbox(slide, 1.0, y, 7, 0.4, kicker,
                              font_size=12, color=self.theme["ink"], font_name=MONO)
            y += 0.5
        if title:
            self._add_textbox(slide, 1.0, y, 7, 0.7, title,
                              font_size=28, color=self.theme["ink"],
                              font_name=SANS_EN if self.style == "B" else SERIF_ZH)
            y += 0.9
        if body:
            self._add_textbox(slide, 1.0, y, 6.5, 3.0, body,
                              font_size=12, color=self.theme["ink"], font_name=SANS_ZH)
        if callout:
            self._add_accent_bar(slide, 1.0, 5.5, 0.04, 0.7)
            self._add_textbox(slide, 1.2, 5.5, 6, 0.6, callout,
                              font_size=12, color=self.theme["ink"],
                              font_name=SERIF_ZH if self.style == "A" else SANS_ZH)
            if callout_src:
                self._add_textbox(slide, 1.2, 6.0, 6, 0.4, callout_src,
                                  font_size=10,
                                  color=self.theme.get("grey_3", "#737373"),
                                  font_name=MONO)

        if image:
            self._add_image(slide, 7.8, 1.2, 4.8, 5.5, image)

    def build_kpi_tower(self, slide_spec: dict):
        """S06: KPI Tower — vertical bars with data."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg_color = self.theme["paper"]
        self._add_bg(slide, bg_color)

        title = slide_spec.get("title", "")
        kicker = slide_spec.get("kicker", "")
        kpis = slide_spec.get("kpis", [])

        y = 0.8
        if kicker:
            self._add_textbox(slide, 1.0, y, 10, 0.4, kicker,
                              font_size=12, color=self.theme["ink"], font_name=MONO)
            y += 0.5
        if title:
            self._add_textbox(slide, 1.0, y, 10, 0.6, title,
                              font_size=26, color=self.theme["ink"],
                              font_name=SANS_EN)
            y += 0.8

        n = len(kpis)
        bar_w = 2.5
        gap = 0.4
        start_x = 1.0
        max_val = max(k.get("value", 0) for k in kpis) if kpis else 100

        for i, kpi in enumerate(kpis):
            bx = start_x + i * (bar_w + gap)
            label = kpi.get("label", "")
            value = kpi.get("value", 0)
            unit = kpi.get("unit", "")
            note = kpi.get("note", "")
            bar_height = (value / max_val) * 3.0 if max_val > 0 else 1.5

            # Bar
            bar_y = y + 3.0 - bar_height
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(bx), Inches(bar_y), Inches(bar_w - 0.1), Inches(bar_height)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = hex_to_rgb(self.theme["accent"])
            bar.line.fill.background()

            self._add_textbox(slide, bx, y + 3.0 + 0.1, bar_w, 0.35,
                              f"{value}{unit}", font_size=16, bold=True,
                              color=self.theme["accent"], font_name=SANS_EN,
                              alignment=PP_ALIGN.CENTER)
            self._add_textbox(slide, bx, y + 3.0 + 0.45, bar_w, 0.35,
                              label, font_size=11, color=self.theme["ink"],
                              font_name=SANS_ZH, alignment=PP_ALIGN.CENTER)
            if note:
                self._add_textbox(slide, bx, y + 3.0 + 0.75, bar_w, 0.3,
                                  note, font_size=9,
                                  color=self.theme.get("grey_3", "#737373"),
                                  font_name=SANS_ZH, alignment=PP_ALIGN.CENTER)

    def build_closing(self, slide_spec: dict):
        """S10 / closing: Thank you / CTA slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        is_dark = slide_spec.get("theme_class", "hero dark").endswith("dark")
        bg_color = self.theme["ink"] if is_dark else self.theme["paper"]
        text_color = self.theme["paper"] if is_dark else self.theme["ink"]
        self._add_bg(slide, bg_color)

        title = slide_spec.get("title", "Thank You")
        lead = slide_spec.get("lead", "")
        meta = slide_spec.get("meta", "")

        self._add_textbox(slide, 2, 2.5, 9, 1.5, title,
                          font_size=48, color=text_color,
                          font_name=SANS_EN if self.style == "B" else SERIF_ZH,
                          alignment=PP_ALIGN.CENTER)
        if lead:
            self._add_textbox(slide, 2, 4.2, 9, 1.0, lead,
                              font_size=16, color=text_color,
                              font_name=SANS_ZH, alignment=PP_ALIGN.CENTER)
        if meta:
            self._add_textbox(slide, 2, 5.5, 9, 0.5, meta,
                              font_size=12, color=text_color,
                              font_name=MONO, alignment=PP_ALIGN.CENTER)

        # Bottom accent bar centered
        self._add_accent_bar(slide, 5.5, 6.5, 2.3, 0.04)

    # ── Master dispatch ──

    BUILDERS = {
        "hero_cover": build_hero_cover,
        "act_divider": build_act_divider,
        "big_numbers": build_big_numbers,
        "quote_image": build_quote_image,
        "image_grid": build_image_grid,
        "pipeline": build_pipeline,
        "hero_question": build_hero_question,
        "big_quote": build_big_quote,
        "compare": build_compare,
        "mixed_text_image": build_mixed_text_image,
        "kpi_tower": build_kpi_tower,
        "closing": build_closing,
    }

    def build(self, output_path: str):
        slides = self.spec.get("slides", [])
        for slide_spec in slides:
            layout = slide_spec.get("layout", "hero_cover")
            builder = self.BUILDERS.get(layout, self.build_hero_cover)
            builder(self, slide_spec)

        self.prs.save(output_path)
        print(f"[OK] PPTX saved: {output_path}")
        print(f"     Slides: {len(slides)} | Style: {self.style} | Theme key: {self.spec.get('theme', 'default')}")


# ═══════════════════════════════════════════════
# CLI
# ═══════════════════════════════════════════════

def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    spec_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else "output.pptx"

    if spec_path == "-":
        spec = json.load(sys.stdin)
    else:
        with open(spec_path, "r", encoding="utf-8") as f:
            spec = json.load(f)

    spec["_base_dir"] = os.path.dirname(os.path.abspath(spec_path)) if spec_path != "-" else os.getcwd()

    builder = PPTXBuilder(spec)
    builder.build(output_path)


if __name__ == "__main__":
    main()
